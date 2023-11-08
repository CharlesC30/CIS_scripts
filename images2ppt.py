import os
from pptx import Presentation
from pptx.util import Inches

# Function to create a PowerPoint presentation from images in a directory
def create_ppt_from_images(image_dir, output_pptx):
    prs = Presentation()

    # Get a list of image files in the directory
    image_files = [f for f in os.listdir(image_dir) if f.endswith(('.jpg', '.png', '.jpeg', '.gif', '.bmp'))]
    image_files.sort()

    for image_file in image_files:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout

        # title = slide.shapes.title
        # title.text = image_file  # Use the image filename as the slide title
        # title.top = Inches(0.5)

        left = Inches(0.5)
        top = Inches(1)
        pic = slide.shapes.add_picture(os.path.join(image_dir, image_file), left, top, width=Inches(8), height=Inches(6))

    prs.save(output_pptx)
    print(f"PowerPoint presentation saved to {output_pptx}")

if __name__ == '__main__':
    image_dir = '/zhome/clarkcs/Pictures/pore_detection/bbox_pore_finding_sobel_1/thresh_images/'  # Replace with the path to your image directory
    output_pptx = '/zhome/clarkcs/Documents/presentations/flawed_pore_detection.pptx'  # Specify the output PowerPoint file

    create_ppt_from_images(image_dir, output_pptx)
